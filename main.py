# import streamlit as st
# from PyPDF2 import PdfReader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain.embeddings.openai import OpenAIEmbeddings
# from langchain.vectorstores import FAISS
# from langchain.chains.question_answering import load_qa_chain
# from langchain_community.chat_models import ChatOpenAI
#
# OPENAI_API_KEY = "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" #Pass your key here
#
# #Upload PDF files
# st.header("My first Chatbot")
#
# with st.sidebar:
#     st.title("Your Documents")
#     file = st.file_uploader(" Upload a PDf file and start asking questions", type="pdf")
#
# #Extract the text
# if file is not None:
#     pdf_reader = PdfReader(file)
#     text = ""
#     for page in pdf_reader.pages:
#         text += page.extract_text()
#         #st.write(text)
#
# #Break it into chunks
#     text_splitter = RecursiveCharacterTextSplitter(
#         separators="\n",
#         chunk_size=1000,
#         chunk_overlap=150,
#         length_function=len
#     )
#     chunks = text_splitter.split_text(text)
#     #st.write(chunks)
#     # generating embedding
#     embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
#
#     # creating vector store - FAISS
#     vector_store = FAISS.from_texts(chunks, embeddings)
#
#     # get user question
#     user_question = st.text_input("Type Your question here")
#
#     # do similarity search
#     if user_question:
#         match = vector_store.similarity_search(user_question)
#         #st.write(match)
#
#         #define the LLM
#         llm = ChatOpenAI(
#             openai_api_key = OPENAI_API_KEY,
#             temperature = 0,
#             max_tokens = 1000,
#             model_name = "gpt-3.5-turbo"
#         )
#
#         #output results
#         #chain -> take the question, get relevant document, pass it to the LLM, generate the output
#         chain = load_qa_chain(llm, chain_type="stuff")
#         response = chain.run(input_documents = match, question = user_question)
#         st.write(response)
# import streamlit as st
# from PyPDF2 import PdfReader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain.embeddings.openai import OpenAIEmbeddings
# from langchain.vectorstores import FAISS
# from langchain.chains.question_answering import load_qa_chain
# from langchain_community.chat_models import ChatOpenAI
# import docx
# from pptx import Presentation
# import pandas as pd
#
# # Initialize LLM
# OPENAI_API_KEY = "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" # Replace with your OpenAI API key
# llm = ChatOpenAI(
#     openai_api_key=OPENAI_API_KEY,
#     temperature=0,
#     max_tokens=1000,
#     model_name="gpt-3.5-turbo"
# )
#
# # Function to handle PDF text extraction
# def extract_text_from_pdf(file):
#     pdf_reader = PdfReader(file)
#     text = ""
#     for page in pdf_reader.pages:
#         text += page.extract_text()
#     return text
#
# # Function to handle DOCX text extraction
# def extract_text_from_docx(file):
#     doc = docx.Document(file)
#     text = ""
#     for para in doc.paragraphs:
#         text += para.text
#     return text
#
# # Function to handle PPTX text extraction
# def extract_text_from_pptx(file):
#     prs = Presentation(file)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text
#     return text
#
# # Function to handle Excel text extraction
# def extract_text_from_excel(file):
#     df = pd.read_excel(file)
#     text = df.to_string(index=False)# Convert DataFrame to string
#     # print(text)
#     return text
#
# # Function to initialize embeddings and vector store
# def initialize_embeddings_and_vector_store(chunks):
#     embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
#     vector_store = FAISS.from_texts(chunks, embeddings)
#     return embeddings, vector_store
#
# # Function to process user question and generate response
# def process_user_question(user_question, vector_store):
#     match = vector_store.similarity_search(user_question)
#     chain = load_qa_chain(llm, chain_type="stuff")
#     response = chain.run(input_documents=match, question=user_question)
#     return response
#
# # Main function for Streamlit app
# def main():
#     st.header("My First Chatbot")
#
#     # Sidebar for file upload
#     with st.sidebar:
#         st.title("Your Documents")
#         file = st.file_uploader("Upload a document (PDF, DOCX, PPTX, or Excel) and start asking questions", type=["pdf", "docx", "pptx", "xlsx"])
#
#     if file is not None:
#         if file.type == "application/pdf":
#             text = extract_text_from_pdf(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
#             text = extract_text_from_docx(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
#             text = extract_text_from_pptx(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
#             text = extract_text_from_excel(file)
#         else:
#             st.error("Unsupported file format")
#
#         text_splitter = RecursiveCharacterTextSplitter(
#             separators="\n",
#             chunk_size=1000,
#             chunk_overlap=150,
#             length_function=len
#         )
#         chunks = text_splitter.split_text(text)
#
#         embeddings, vector_store = initialize_embeddings_and_vector_store(chunks)
#
#         user_question = st.text_input("Type Your question here")
#
#         if user_question:
#             response = process_user_question(user_question, vector_store)
#             st.write(response)
#
# if __name__ == "__main__":
#     main()
# import streamlit as st
# from PyPDF2 import PdfReader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain.embeddings.openai import OpenAIEmbeddings
# from langchain.vectorstores import FAISS
# from langchain.chains.question_answering import load_qa_chain
# from langchain_community.chat_models import ChatOpenAI
# import docx
# from pptx import Presentation
# import pandas as pd
#
# # Initialize LLM
# OPENAI_API_KEY = "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" # Replace with your OpenAI API key
# llm = ChatOpenAI(
#     openai_api_key=OPENAI_API_KEY,
#     temperature=0,
#     max_tokens=1000,
#     model_name="gpt-3.5-turbo"
# )
#
# # Function to handle PDF text extraction
# def extract_text_from_pdf(file):
#     pdf_reader = PdfReader(file)
#     text = ""
#     for page in pdf_reader.pages:
#         text += page.extract_text()
#     return text
#
# # Function to handle DOCX text extraction
# def extract_text_from_docx(file):
#     doc = docx.Document(file)
#     text = ""
#     for para in doc.paragraphs:
#         text += para.text
#     return text
#
# # Function to handle PPTX text extraction
# def extract_text_from_pptx(file):
#     prs = Presentation(file)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text
#     return text
#
# # Function to handle Excel text extraction
# def extract_text_from_excel(file):
#     df = pd.read_excel(file)
#     text = df.to_string(index=False)# Convert DataFrame to string
#     return text
#
# # Function to initialize embeddings and vector store
# def initialize_embeddings_and_vector_store(chunks):
#     embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
#     vector_store = FAISS.from_texts(chunks, embeddings)
#     return embeddings, vector_store
#
# # Function to process user question and generate response
# def process_user_question(user_question, vector_store):
#     match = vector_store.similarity_search(user_question)
#     chain = load_qa_chain(llm, chain_type="stuff")
#     response = chain.run(input_documents=match, question=user_question)
#     return response
#
# # Main function for Streamlit app
# def main():
#     st.title("Chatbot")
#
#     # Sidebar for file upload
#     with st.sidebar:
#         st.title("Your Documents")
#         file = st.file_uploader("Upload a document (PDF, DOCX, PPTX, or Excel) and start asking questions", type=["pdf", "docx", "pptx", "xlsx"])
#
#     if file is not None:
#         if file.type == "application/pdf":
#             text = extract_text_from_pdf(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
#             text = extract_text_from_docx(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
#             text = extract_text_from_pptx(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
#             text = extract_text_from_excel(file)
#         else:
#             st.error("Unsupported file format")
#
#         text_splitter = RecursiveCharacterTextSplitter(
#             separators="\n",
#             chunk_size=1000,
#             chunk_overlap=150,
#             length_function=len
#         )
#         chunks = text_splitter.split_text(text)
#
#         embeddings, vector_store = initialize_embeddings_and_vector_store(chunks)
#
#         chat_history = []
#
#         while True:
#             user_question = st.text_input("You:", key=f"user_input_{len(chat_history)}")
#             if user_question:
#                 chat_history.append({"user": user_question})
#                 response = process_user_question(user_question, vector_store)
#                 chat_history.append({"bot": response})
#             else:
#                 break
#
#         for entry in chat_history:
#             if "user" in entry:
#                 st.text("You: " + entry["user"])
#             elif "bot" in entry:
#                 st.text("Bot: " + entry["bot"])
#
# if __name__ == "__main__":
#     main()
# import streamlit as st
# from PyPDF2 import PdfReader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain.embeddings.openai import OpenAIEmbeddings
# from langchain.vectorstores import FAISS
# from langchain.chains.question_answering import load_qa_chain
# from langchain_community.chat_models import ChatOpenAI
# import docx
# from pptx import Presentation
# import pandas as pd
#
# # Initialize LLM
# OPENAI_API_KEY = "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" # Replace with your OpenAI API key
# llm = ChatOpenAI(
#     openai_api_key=OPENAI_API_KEY,
#     temperature=0,
#     max_tokens=1000,
#     model_name="gpt-3.5-turbo"
# )
#
# # Function to handle PDF text extraction
# def extract_text_from_pdf(file):
#     pdf_reader = PdfReader(file)
#     text = ""
#     for page in pdf_reader.pages:
#         text += page.extract_text()
#     return text
#
# # Function to handle DOCX text extraction
# def extract_text_from_docx(file):
#     doc = docx.Document(file)
#     text = ""
#     for para in doc.paragraphs:
#         text += para.text
#     return text
#
# # Function to handle PPTX text extraction
# def extract_text_from_pptx(file):
#     prs = Presentation(file)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text
#     return text
#
# # Function to handle Excel text extraction
# def extract_text_from_excel(file):
#     df = pd.read_excel(file)
#     text = df.to_string(index=False)  # Convert DataFrame to string
#     return text
#
# # Function to initialize embeddings and vector store
# def initialize_embeddings_and_vector_store(chunks):
#     embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
#     vector_store = FAISS.from_texts(chunks, embeddings)
#     return embeddings, vector_store
#
# # Function to process user question and generate response
# def process_user_question(user_question, vector_store):
#     match = vector_store.similarity_search(user_question)
#     chain = load_qa_chain(llm, chain_type="stuff")
#     response = chain.run(input_documents=match, question=user_question)
#     return response
#
# # Main function for Streamlit app
# def main():
#     st.title("Chatbot")
#
#     # Sidebar for file upload
#     with st.sidebar:
#         st.title("Your Documents")
#         file = st.file_uploader("Upload a document (PDF, DOCX, PPTX, or Excel) and start asking questions",
#                                 type=["pdf", "docx", "pptx", "xlsx"])
#
#     if file is not None:
#         if file.type == "application/pdf":
#             text = extract_text_from_pdf(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
#             text = extract_text_from_docx(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
#             text = extract_text_from_pptx(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
#             text = extract_text_from_excel(file)
#         else:
#             st.error("Unsupported file format")
#
#         text_splitter = RecursiveCharacterTextSplitter(
#             separators="\n",
#             chunk_size=1000,
#             chunk_overlap=150,
#             length_function=len
#         )
#         chunks = text_splitter.split_text(text)
#
#         embeddings, vector_store = initialize_embeddings_and_vector_store(chunks)
#
#         question_inputs = []
#
#         while True:
#             question_index = len(question_inputs)
#             user_question = st.text_input(f"Question {question_index + 1}:", key=f'user_input_{question_index}')
#             if user_question == "":
#                 break
#             if user_question:
#                 question_inputs.append(user_question)
#                 response = process_user_question(user_question, vector_store)
#                 st.text("Bot: " + response)
#
# if __name__ == "__main__":
#     main()
#
# import streamlit as st
# from PyPDF2 import PdfReader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain.embeddings.openai import OpenAIEmbeddings
# from langchain.vectorstores import FAISS
# from langchain.chains.question_answering import load_qa_chain
# from langchain_community.chat_models import ChatOpenAI
# import docx
# from pptx import Presentation
# import pandas as pd
#
# # Initialize LLM
# OPENAI_API_KEY = "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" # Replace with your OpenAI API key
# llm = ChatOpenAI(
#     openai_api_key=OPENAI_API_KEY,
#     temperature=0,
#     max_tokens=1000,
#     model_name="gpt-3.5-turbo"
# )
#
# # Function to handle PDF text extraction
# def extract_text_from_pdf(file):
#     pdf_reader = PdfReader(file)
#     text = ""
#     for page in pdf_reader.pages:
#         text += page.extract_text()
#     return text
#
# # Function to handle DOCX text extraction
# def extract_text_from_docx(file):
#     doc = docx.Document(file)
#     text = ""
#     for para in doc.paragraphs:
#         text += para.text
#     return text
#
# # Function to handle PPTX text extraction
# def extract_text_from_pptx(file):
#     prs = Presentation(file)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text
#     return text
#
# # Function to handle Excel text extraction
# def extract_text_from_excel(file):
#     df = pd.read_excel(file)
#     text = df.to_string(index=False)  # Convert DataFrame to string
#     return text
#
# # Function to initialize embeddings and vector store
# def initialize_embeddings_and_vector_store(chunks):
#     embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
#     vector_store = FAISS.from_texts(chunks, embeddings)
#     return embeddings, vector_store
#
# # Function to process user question and generate response
# def process_user_question(user_question, vector_store):
#     match = vector_store.similarity_search(user_question)
#     chain = load_qa_chain(llm, chain_type="stuff")
#     response = chain.run(input_documents=match, question=user_question)
#     return response
#
# # Main function for Streamlit app
# def main():
#     st.title("Chatbot")
#
#     # Sidebar for file upload
#     with st.sidebar:
#         st.title("Your Documents")
#         file = st.file_uploader("Upload a document (PDF, DOCX, PPTX, or Excel) and start asking questions",
#                                 type=["pdf", "docx", "pptx", "xlsx"])
#
#     if file is not None:
#         if file.type == "application/pdf":
#             text = extract_text_from_pdf(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
#             text = extract_text_from_docx(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
#             text = extract_text_from_pptx(file)
#         elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
#             text = extract_text_from_excel(file)
#         else:
#             st.error("Unsupported file format")
#
#         text_splitter = RecursiveCharacterTextSplitter(
#             separators="\n",
#             chunk_size=1000,
#             chunk_overlap=150,
#             length_function=len
#         )
#         chunks = text_splitter.split_text(text)
#
#         embeddings, vector_store = initialize_embeddings_and_vector_store(chunks)
#
#         question_inputs = []
#
#         while True:
#             question_index = len(question_inputs)
#             user_question = st.text_input(f"Question {question_index + 1}:", key=f'user_input_{question_index}')
#             if user_question == "":
#                 break
#             if user_question:
#                 question_inputs.append(user_question)
#                 response = process_user_question(user_question, vector_store)
#                 st.write(f"You: {user_question}")
#                 st.write("Bot:", response)
#
# if __name__ == "__main__":
#     main()
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain_community.chat_models import ChatOpenAI
import docx
from pptx import Presentation
import pandas as pd

# Initialize LLM
OPENAI_API_KEY = "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" # Replace with your OpenAI API key
llm = ChatOpenAI(
    openai_api_key=OPENAI_API_KEY,
    temperature=0,
    max_tokens=1000,
    model_name="gpt-3.5-turbo"
)

# Function to handle PDF text extraction
def extract_text_from_pdf(file):
    pdf_reader = PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

# Function to handle DOCX text extraction
def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text

# Function to handle PPTX text extraction
def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

# Function to handle Excel text extraction
def extract_text_from_excel(file):
    df = pd.read_excel(file)
    text = df.to_string(index=False)  # Convert DataFrame to string
    return text

# Function to initialize embeddings and vector store
def initialize_embeddings_and_vector_store(chunks):
    embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
    vector_store = FAISS.from_texts(chunks, embeddings)
    return embeddings, vector_store

# Function to process user question and generate response
def process_user_question(user_question, vector_store):
    match = vector_store.similarity_search(user_question)
    chain = load_qa_chain(llm, chain_type="stuff")
    response = chain.run(input_documents=match, question=user_question)
    return response

# Main function for Streamlit app
def main():
    st.title("Chatbot")

    # Sidebar for file upload
    st.sidebar.title("Your Documents")
    files = st.sidebar.file_uploader("Upload one or more documents (PDF, DOCX, PPTX, or Excel) to start asking questions",
                                      accept_multiple_files=True)

    if files:
        text = ""
        for file in files:
            if file.type == "application/pdf":
                text += extract_text_from_pdf(file)
            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text += extract_text_from_docx(file)
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text += extract_text_from_pptx(file)
            elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                text += extract_text_from_excel(file)
            else:
                st.error(f"Unsupported file format: {file.name}")

        text_splitter = RecursiveCharacterTextSplitter(
            separators="\n",
            chunk_size=1000,
            chunk_overlap=150,
            length_function=len
        )
        chunks = text_splitter.split_text(text)

        embeddings, vector_store = initialize_embeddings_and_vector_store(chunks)

        question_inputs = []

        while True:
            question_index = len(question_inputs)
            user_question = st.text_input(f"Question {question_index + 1}:", key=f'user_input_{question_index}')
            if user_question == "":
                break
            if user_question:
                question_inputs.append(user_question)
                response = process_user_question(user_question, vector_store)
                st.write(f"You: {user_question}")
                st.write("Bot:", response)

if __name__ == "__main__":
    main()




